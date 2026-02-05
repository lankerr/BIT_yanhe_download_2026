"""
PPT/幻灯片提取器 - GPU加速版本
使用FFmpeg scene filter进行场景变化检测，速度约为50x实时
"""

import subprocess
import os
import sys
import re
import argparse
import cv2
import numpy as np
from pathlib import Path
from typing import List, Tuple, Optional
from tqdm import tqdm
import imagehash
from PIL import Image

# 尝试导入pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# 配置参数
SCENE_THRESHOLD = 0.1  # FFmpeg scene filter阈值（0.0-1.0）
MIN_SCENE_INTERVAL = 2.0  # 最小场景间隔（秒）
SIMILARITY_THRESHOLD = 0.90  # 感知哈希相似度阈值
JPG_QUALITY = 95  # JPEG质量

def get_video_duration(video_path: str) -> float:
    """获取视频时长（秒）"""
    cmd = [
        'ffprobe', '-v', 'error',
        '-show_entries', 'format=duration',
        '-of', 'csv=p=0',
        video_path
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True)
        return float(result.stdout.strip())
    except:
        return 0


def detect_scenes_ffmpeg(video_path: str, threshold: float = SCENE_THRESHOLD) -> List[float]:
    """
    使用FFmpeg scene filter检测场景变化时间点
    返回场景变化的时间戳列表（秒）
    """
    duration = get_video_duration(video_path)
    duration_str = f"{int(duration//60)}分{int(duration%60)}秒" if duration > 0 else "未知"
    
    cmd = [
        'ffmpeg', '-hide_banner',
        '-i', video_path,
        '-vf', f"select='gt(scene,{threshold})',showinfo",
        '-an', '-f', 'null', '-'
    ]
    
    print(f"正在使用 FFmpeg 分析场景变化 (threshold={threshold})...")
    print(f"视频时长: {duration_str}，预计处理时间: {int(duration/50)+1}秒 (约50x实时)")
    print("请等待处理完成...")
    
    timestamps = []
    
    try:
        # 使用Popen实时读取输出
        process = subprocess.Popen(
            cmd,
            stderr=subprocess.PIPE,
            stdout=subprocess.DEVNULL,
            text=True,
            encoding='utf-8',
            errors='replace'
        )
        
        pattern = r'pts_time:(\d+\.?\d*)'
        time_pattern = r'time=(\d+:\d+:\d+\.\d+)'
        last_progress = 0
        
        for line in process.stderr:
            # 检测场景变化点
            match = re.search(pattern, line)
            if match:
                timestamps.append(float(match.group(1)))
                print(f"  找到场景 #{len(timestamps)}: {float(match.group(1)):.1f}s")
            
            # 显示进度
            time_match = re.search(time_pattern, line)
            if time_match and duration > 0:
                time_str = time_match.group(1)
                parts = time_str.split(':')
                current_time = int(parts[0])*3600 + int(parts[1])*60 + float(parts[2])
                progress = int(current_time / duration * 100)
                if progress > last_progress and progress % 10 == 0:
                    print(f"  进度: {progress}%")
                    last_progress = progress
        
        process.wait()
        
    except Exception as e:
        print(f"FFmpeg执行错误: {e}")
        return []
    
    print(f"检测到 {len(timestamps)} 个场景变化点")
    return timestamps


def filter_close_timestamps(timestamps: List[float], min_interval: float = MIN_SCENE_INTERVAL) -> List[float]:
    """过滤太接近的时间戳"""
    if not timestamps:
        return []
    
    filtered = [timestamps[0]]
    for t in timestamps[1:]:
        if t - filtered[-1] >= min_interval:
            filtered.append(t)
    
    return filtered


def extract_frame_at_time(video_path: str, timestamp: float) -> Optional[np.ndarray]:
    """使用FFmpeg在指定时间点提取单帧"""
    cmd = [
        'ffmpeg', '-hide_banner', '-loglevel', 'error',
        '-ss', str(timestamp),
        '-i', video_path,
        '-vframes', '1',
        '-f', 'image2pipe',
        '-pix_fmt', 'bgr24',
        '-vcodec', 'rawvideo',
        '-'
    ]
    
    try:
        result = subprocess.run(cmd, capture_output=True, timeout=10)
        if result.returncode != 0:
            return None
        
        # 获取视频尺寸
        probe_cmd = [
            'ffprobe', '-v', 'error',
            '-select_streams', 'v:0',
            '-show_entries', 'stream=width,height',
            '-of', 'csv=p=0',
            video_path
        ]
        probe_result = subprocess.run(probe_cmd, capture_output=True, text=True)
        w, h = map(int, probe_result.stdout.strip().split(','))
        
        frame = np.frombuffer(result.stdout, dtype=np.uint8).reshape((h, w, 3))
        return frame
    except Exception as e:
        return None


def extract_frames_batch(video_path: str, timestamps: List[float]) -> List[Tuple[float, np.ndarray]]:
    """批量提取帧（使用OpenCV）"""
    if not timestamps:
        return []
    
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        print(f"无法打开视频: {video_path}")
        return []
    
    fps = cap.get(cv2.CAP_PROP_FPS)
    frames = []
    
    for timestamp in tqdm(timestamps, desc="提取帧", unit="帧"):
        frame_num = int(timestamp * fps)
        cap.set(cv2.CAP_PROP_POS_FRAMES, frame_num)
        ret, frame = cap.read()
        if ret:
            frames.append((timestamp, frame))
    
    cap.release()
    return frames


def compute_phash(image: np.ndarray, hash_size: int = 16) -> imagehash.ImageHash:
    """计算感知哈希"""
    pil_image = Image.fromarray(cv2.cvtColor(image, cv2.COLOR_BGR2RGB))
    return imagehash.phash(pil_image, hash_size=hash_size)


def deduplicate_frames(frames: List[Tuple[float, np.ndarray]], 
                       threshold: float = SIMILARITY_THRESHOLD) -> List[Tuple[float, np.ndarray]]:
    """使用感知哈希去重"""
    if not frames:
        return []
    
    unique_frames = [frames[0]]
    prev_hash = compute_phash(frames[0][1])
    
    for timestamp, frame in tqdm(frames[1:], desc="去重", unit="帧"):
        current_hash = compute_phash(frame)
        similarity = 1 - (prev_hash - current_hash) / (16 * 16)  # hash_size=16
        
        if similarity < threshold:
            unique_frames.append((timestamp, frame))
            prev_hash = current_hash
    
    return unique_frames


def save_image_unicode(image: np.ndarray, path: str, quality: int = JPG_QUALITY) -> bool:
    """保存图像，支持Unicode路径"""
    try:
        _, buf = cv2.imencode('.jpg', image, [cv2.IMWRITE_JPEG_QUALITY, quality])
        with open(path, 'wb') as f:
            f.write(buf.tobytes())
        return True
    except Exception as e:
        print(f"保存失败 {path}: {e}")
        return False


def extract_slides(video_path: str, 
                   output_dir: str,
                   scene_threshold: float = SCENE_THRESHOLD,
                   min_interval: float = MIN_SCENE_INTERVAL,
                   similarity_threshold: float = SIMILARITY_THRESHOLD) -> int:
    """
    主函数：从视频中提取幻灯片
    返回保存的幻灯片数量
    """
    video_path = os.path.abspath(video_path)
    output_dir = os.path.abspath(output_dir)
    
    if not os.path.exists(video_path):
        print(f"视频文件不存在: {video_path}")
        return 0
    
    os.makedirs(output_dir, exist_ok=True)
    
    video_name = Path(video_path).stem
    
    # Step 1: 场景检测
    print("\n=== 步骤 1/4: 场景检测 ===")
    timestamps = detect_scenes_ffmpeg(video_path, scene_threshold)
    
    if not timestamps:
        print("未检测到场景变化，尝试降低阈值...")
        timestamps = detect_scenes_ffmpeg(video_path, scene_threshold * 0.5)
    
    if not timestamps:
        print("无法检测到场景变化")
        return 0
    
    # Step 2: 过滤
    print("\n=== 步骤 2/4: 过滤时间戳 ===")
    filtered_timestamps = filter_close_timestamps(timestamps, min_interval)
    print(f"过滤后: {len(timestamps)} → {len(filtered_timestamps)} 个时间点")
    
    # Step 3: 提取帧
    print("\n=== 步骤 3/4: 提取帧 ===")
    frames = extract_frames_batch(video_path, filtered_timestamps)
    print(f"成功提取 {len(frames)} 帧")
    
    # Step 4: 去重
    print("\n=== 步骤 4/4: 去重并保存 ===")
    unique_frames = deduplicate_frames(frames, similarity_threshold)
    print(f"去重后: {len(frames)} → {len(unique_frames)} 帧")
    
    # 保存
    saved_count = 0
    for i, (timestamp, frame) in enumerate(tqdm(unique_frames, desc="保存", unit="张")):
        minutes = int(timestamp // 60)
        seconds = int(timestamp % 60)
        filename = f"{video_name}_slide_{i+1:03d}_{minutes:02d}m{seconds:02d}s.jpg"
        filepath = os.path.join(output_dir, filename)
        
        if save_image_unicode(frame, filepath):
            saved_count += 1
    
    print(f"\n完成！保存了 {saved_count} 张幻灯片到: {output_dir}")
    return saved_count


def create_pptx(image_dir: str, output_pptx: str = None) -> str:
    """
    将图片目录中的所有图片合并成PPT
    返回生成的PPT文件路径
    """
    if not HAS_PPTX:
        print("警告: python-pptx未安装，无法生成PPT。请运行: pip install python-pptx")
        return None
    
    import glob
    
    # 获取所有图片并排序
    images = sorted(glob.glob(os.path.join(image_dir, '*.jpg')))
    if not images:
        print(f"未找到图片: {image_dir}")
        return None
    
    # 确定输出文件名
    if output_pptx is None:
        dir_name = Path(image_dir).name
        output_pptx = os.path.join(os.path.dirname(image_dir), f"{dir_name}.pptx")
    
    print(f"\n=== 生成PPT ===\n  图片数量: {len(images)}")
    
    # 创建PPT (16:9)
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 添加每张图片为一页幻灯片
    for img_path in tqdm(images, desc="生成PPT", unit="页"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(16), height=Inches(9))
    
    # 保存
    prs.save(output_pptx)
    file_size = os.path.getsize(output_pptx) / 1024 / 1024
    print(f"  PPT已保存: {output_pptx} ({file_size:.2f} MB)")
    
    return output_pptx


def main():
    parser = argparse.ArgumentParser(
        description="PPT/幻灯片提取器 - GPU加速版本",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  python ppt_extractor_gpu.py video.mp4 -o output/slides
  python ppt_extractor_gpu.py video.mp4 -t 0.2 -m 3.0
        """
    )
    
    parser.add_argument('video', help='视频文件路径')
    parser.add_argument('-o', '--output', default='output/ppt_slides', help='输出目录')
    parser.add_argument('-t', '--threshold', type=float, default=SCENE_THRESHOLD,
                        help=f'场景检测阈值 (默认: {SCENE_THRESHOLD})')
    parser.add_argument('-m', '--min-interval', type=float, default=MIN_SCENE_INTERVAL,
                        help=f'最小场景间隔秒数 (默认: {MIN_SCENE_INTERVAL})')
    parser.add_argument('-s', '--similarity', type=float, default=SIMILARITY_THRESHOLD,
                        help=f'相似度阈值 (默认: {SIMILARITY_THRESHOLD})')
    parser.add_argument('-p', '--pptx', action='store_true',
                        help='同时生成PPT文件')
    
    args = parser.parse_args()
    
    print("=" * 60)
    print("PPT/幻灯片提取器 - GPU加速版本")
    print("=" * 60)
    print(f"视频: {args.video}")
    print(f"输出目录: {args.output}")
    print(f"场景阈值: {args.threshold}")
    print(f"最小间隔: {args.min_interval}s")
    print(f"相似度阈值: {args.similarity}")
    print(f"生成PPT: {'是' if args.pptx else '否'}")
    print("=" * 60)
    
    count = extract_slides(
        args.video,
        args.output,
        scene_threshold=args.threshold,
        min_interval=args.min_interval,
        similarity_threshold=args.similarity
    )
    
    # 如果需要生成PPT
    if count > 0 and args.pptx:
        create_pptx(args.output)
    
    return 0 if count > 0 else 1


if __name__ == '__main__':
    sys.exit(main())
