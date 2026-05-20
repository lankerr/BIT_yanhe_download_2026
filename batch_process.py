"""
批量后处理器 - 统一处理下载的课程视频
功能：
1. GPU 音频转录（faster-whisper）
2. GPU 课件提取（FFmpeg scene detection）
3. 可单独或组合使用

加固点（2026-04 补充）：
- collect_videos 过滤损坏 mp4（<1MB 或时长<30s）, 防止流水线迫件
- batch_extract_slides 在 PPT 提取失败时生成占位 pptx，避免 watchdog 死循环
- batch_extract_slides / batch_transcribe_videos 单视频崩溃雔离, 不拖败整个课程
"""
import os
import sys
import time
import subprocess
import argparse
from pathlib import Path
from typing import Optional, Callable, List

# 强制 stdout/stderr 为 UTF-8 (Windows 后台运行时默认 GBK 会因 emoji 崩溃)
try:
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
    sys.stderr.reconfigure(encoding='utf-8', errors='replace')
except Exception:
    pass

# 健康检查阈值
MIN_VIDEO_SIZE_BYTES = 1 * 1024 * 1024   # 1 MB
MIN_VIDEO_DURATION_SEC = 30              # 30 秒


try:
    from app_paths import ffprobe_path
except ImportError:
    def ffprobe_path(): return "ffprobe"


def _ffprobe_duration(path: str) -> float:
    """返回视频时长(秒)，失败返 0"""
    try:
        out = subprocess.run(
            [ffprobe_path(), "-v", "error", "-show_entries", "format=duration",
             "-of", "default=noprint_wrappers=1:nokey=1", path],
            capture_output=True, text=True, timeout=15,
        ).stdout.strip()
        return float(out) if out else 0.0
    except Exception:
        return 0.0


def _is_video_healthy(path: str) -> bool:
    """尺寸过小或时长过短视为损坏。"""
    try:
        if os.path.getsize(path) < MIN_VIDEO_SIZE_BYTES:
            return False
        if _ffprobe_duration(path) < MIN_VIDEO_DURATION_SEC:
            return False
        return True
    except Exception:
        return False


def _make_placeholder_pptx(pptx_path: str, video_name: str, reason: str) -> bool:
    """生成单页占位 pptx，表示 PPT 提取已尝试但不可行。"""
    try:
        from pptx import Presentation
        os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = video_name
        slide.placeholders[1].text = f"PPT 提取失败: {reason}\n请查看转录文本。"
        prs.save(pptx_path)
        print(f"  已生成占位 pptx: {pptx_path}")
        return True
    except Exception as e:
        print(f"  占位 pptx 生成失败: {e}")
        return False


def collect_videos(input_dir: str) -> List[str]:
    """收集目录中所有 mp4 视频文件，过滤损坏者"""
    mp4_files = []
    skipped_bad = []
    for root, dirs, files in os.walk(input_dir):
        if any(skip in root for skip in ["extracted_ppt", "transcripts", "extracted_slides", "_broken"]):
            continue
        for f in sorted(files):
            if f.lower().endswith(".mp4"):
                full = os.path.join(root, f)
                if _is_video_healthy(full):
                    mp4_files.append(full)
                else:
                    skipped_bad.append(full)
    if skipped_bad:
        print(f"⚠️  跳过 {len(skipped_bad)} 个损坏/过短视频:")
        for p in skipped_bad:
            print(f"     - {os.path.basename(p)}")
    return mp4_files


def batch_extract_slides(
    input_dir: str,
    output_dir: str = None,
    scene_threshold: float = 0.1,
    min_interval: float = 2.0,
    similarity_threshold: float = 0.90,
    progress_callback: Optional[Callable] = None,
) -> List[dict]:
    """
    批量提取课件幻灯片
    """
    import ppt_extractor_gpu as extractor

    if output_dir is None:
        output_dir = os.path.join(input_dir, "extracted_ppt")
    os.makedirs(output_dir, exist_ok=True)

    mp4_files = collect_videos(input_dir)
    if not mp4_files:
        print("未找到视频文件")
        return []

    print(f"{'=' * 60}")
    print(f"批量课件提取")
    print(f"共找到 {len(mp4_files)} 个视频文件")
    print(f"{'=' * 60}")

    results = []
    total_slides = 0

    for i, video_path in enumerate(mp4_files):
        video_name = Path(video_path).stem
        slide_dir = os.path.join(output_dir, video_name)
        pptx_path = os.path.join(output_dir, f"{video_name}.pptx")

        print(f"\n{'=' * 60}")
        print(f"[{i + 1}/{len(mp4_files)}] {video_name}")
        print(f"{'=' * 60}")

        if os.path.exists(pptx_path):
            print(f"  已存在 PPT，跳过")
            results.append({"name": video_name, "status": "跳过", "slides": 0, "elapsed": 0})
            continue

        start = time.time()
        try:
            if progress_callback:
                overall = int(i / len(mp4_files) * 100)
                progress_callback(overall, 100, f"[{i+1}/{len(mp4_files)}] 提取课件: {video_name[:30]}...")

            count = extractor.extract_slides(
                video_path, slide_dir,
                scene_threshold=scene_threshold,
                min_interval=min_interval,
                similarity_threshold=similarity_threshold,
            )
            if count > 0:
                extractor.create_pptx(slide_dir, pptx_path)
                total_slides += count
                status = f"成功({count}张)"
            else:
                # 无场景变化 -> 生成占位 pptx, 避免下次重复处理
                _make_placeholder_pptx(pptx_path, video_name, "无场景变化")
                status = "占位(无场景)"

            elapsed = time.time() - start
            results.append({"name": video_name, "status": status, "slides": count, "elapsed": elapsed})

        except Exception as e:
            elapsed = time.time() - start
            print(f"  失败: {e}")
            # 崩溃同样生成占位, 避免 watchdog 死循环
            _make_placeholder_pptx(pptx_path, video_name, f"处理异常: {type(e).__name__}")
            results.append({"name": video_name, "status": f"失败占位: {e}", "slides": 0, "elapsed": elapsed})

    print(f"\n{'=' * 60}")
    print(f"课件提取完成! 总计: {total_slides} 张")
    print(f"{'=' * 60}")
    for r in results:
        t = f" ({r['elapsed']:.0f}s)" if r['elapsed'] > 0 else ""
        print(f"  {r['name']}: {r['status']}{t}")

    return results


def batch_transcribe_videos(
    input_dir: str,
    output_dir: str = None,
    model_size: str = "large-v3",
    device: str = "auto",
    compute_type: str = "float16",
    language: str = "zh",
    progress_callback: Optional[Callable] = None,
) -> List[dict]:
    """
    批量转录视频音频
    """
    import audio_transcriber_gpu as transcriber_mod

    if output_dir is None:
        output_dir = os.path.join(input_dir, "transcripts")
    os.makedirs(output_dir, exist_ok=True)

    mp4_files = collect_videos(input_dir)
    if not mp4_files:
        print("未找到视频文件")
        return []

    print(f"{'=' * 60}")
    print(f"批量音频转录 (GPU加速)")
    print(f"共找到 {len(mp4_files)} 个视频文件")
    print(f"模型: {model_size}, 设备: {device}")
    print(f"{'=' * 60}")

    transcriber = transcriber_mod.AudioTranscriber(
        model_size=model_size,
        device=device,
        compute_type=compute_type,
        language=language,
    )
    transcriber.load_model(progress_callback)

    results = []
    for i, video_path in enumerate(mp4_files):
        video_name = Path(video_path).stem
        file_output_dir = os.path.join(output_dir, video_name)

        print(f"\n{'=' * 60}")
        print(f"[{i + 1}/{len(mp4_files)}] {video_name}")
        print(f"{'=' * 60}")

        txt_path = os.path.join(file_output_dir, f"{video_name}.txt")
        if os.path.exists(txt_path):
            print(f"  已存在转录文件，跳过")
            results.append({"name": video_name, "status": "跳过", "elapsed": 0})
            continue

        start = time.time()
        try:
            def file_progress(current, total, msg,
                              idx=i, count=len(mp4_files)):
                if progress_callback:
                    overall = int((idx / count + current / total / count) * 100)
                    progress_callback(overall, 100, f"[{idx+1}/{count}] {msg}")

            result = transcriber.transcribe_file(
                video_path,
                output_dir=file_output_dir,
                progress_callback=file_progress,
            )
            elapsed = time.time() - start
            seg_count = len(result.get("segments", []))
            results.append({"name": video_name, "status": f"成功({seg_count}段)", "elapsed": elapsed})

        except Exception as e:
            elapsed = time.time() - start
            print(f"  失败: {e}")
            results.append({"name": video_name, "status": f"失败: {e}", "elapsed": elapsed})

    print(f"\n{'=' * 60}")
    print(f"批量转录完成!")
    print(f"{'=' * 60}")
    for r in results:
        t = f" ({r['elapsed']:.0f}s)" if r['elapsed'] > 0 else ""
        print(f"  {r['name']}: {r['status']}{t}")

    return results


def batch_process_all(
    input_dir: str,
    do_slides: bool = True,
    do_transcribe: bool = True,
    model_size: str = "large-v3",
    device: str = "auto",
    compute_type: str = "float16",
    language: str = "zh",
    progress_callback: Optional[Callable] = None,
) -> dict:
    """
    统一批处理：课件提取 + 音频转录
    """
    all_results = {"slides": [], "transcripts": []}

    if do_slides:
        print("\n" + "=" * 60)
        print("阶段 1: 课件提取")
        print("=" * 60)
        all_results["slides"] = batch_extract_slides(
            input_dir, progress_callback=progress_callback,
        )

    if do_transcribe:
        print("\n" + "=" * 60)
        print("阶段 2: 音频转录")
        print("=" * 60)
        all_results["transcripts"] = batch_transcribe_videos(
            input_dir, model_size=model_size, device=device,
            compute_type=compute_type, language=language,
            progress_callback=progress_callback,
        )

    return all_results


def main():
    parser = argparse.ArgumentParser(
        description="批量后处理器 - 课件提取 + 音频转录",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  python batch_process.py output/                      # 执行全部处理
  python batch_process.py output/ --slides-only        # 仅提取课件
  python batch_process.py output/ --transcribe-only    # 仅转录音频
  python batch_process.py output/ -m medium -l en      # 用 medium 模型转录英文
        """,
    )

    parser.add_argument("input_dir", help="视频所在目录")
    parser.add_argument("--slides-only", action="store_true", help="仅提取课件")
    parser.add_argument("--transcribe-only", action="store_true", help="仅转录音频")
    parser.add_argument("-m", "--model", default="large-v3", help="Whisper 模型 (默认 large-v3)")
    parser.add_argument("--device", default="auto", help="设备 (auto/cuda/cpu)")
    parser.add_argument("--compute-type", default="float16", help="计算类型")
    parser.add_argument("-l", "--language", default="zh", help="语言 (默认 zh)")

    args = parser.parse_args()

    do_slides = not args.transcribe_only
    do_transcribe = not args.slides_only

    batch_process_all(
        args.input_dir,
        do_slides=do_slides,
        do_transcribe=do_transcribe,
        model_size=args.model,
        device=args.device,
        compute_type=args.compute_type,
        language=args.language,
    )


if __name__ == "__main__":
    main()
