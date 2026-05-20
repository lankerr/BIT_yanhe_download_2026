"""破除 watchdog 死循环:
1. 把损坏 mp4 移到 output/_broken/ (watchdog 看不到就不再处理)
2. 对 PPT 提取失败但已转录完成的视频, 生成占位 pptx 满足完成条件
"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

BASE = Path(__file__).parent / "output"
BROKEN = BASE / "_broken"
BROKEN.mkdir(exist_ok=True)

BAD = [
    # (课程目录, 视频文件名 (无后缀), 操作)
    ("科学与工程计算-screen", "科学与工程计算-熊春光-第4周 星期二 第3大节", "move"),
    ("大模型技术原理与实践-screen", "大模型技术原理与实践-林知微-第5周 星期五 第3大节", "placeholder"),
    ("大模型技术原理与实践-screen", "大模型技术原理与实践-林知微-第6周 星期一 第2大节", "placeholder"),
]

def make_placeholder_pptx(path: Path, title: str, note: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = note
    prs.save(str(path))

for course, name, op in BAD:
    course_dir = BASE / course
    mp4 = course_dir / f"{name}.mp4"
    pptx = course_dir / "extracted_ppt" / f"{name}.pptx"
    if op == "move":
        if mp4.exists():
            dst = BROKEN / f"{course}__{name}.mp4"
            shutil.move(str(mp4), str(dst))
            print(f"[MOVED] {mp4.name} -> _broken/")
        else:
            print(f"[SKIP] {mp4.name} 不存在")
    elif op == "placeholder":
        if not pptx.exists():
            pptx.parent.mkdir(parents=True, exist_ok=True)
            make_placeholder_pptx(
                pptx,
                title=name,
                note="PPT 提取失败 (视频内容缺乏场景变化)。请直接查看转录文本。",
            )
            print(f"[PPTX] 占位生成: {pptx.name}")
        else:
            print(f"[SKIP] {pptx.name} 已存在")

print("DONE")
